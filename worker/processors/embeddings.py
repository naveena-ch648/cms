"""Embeddings processor — chunks documents and generates vector embeddings, stores in Qdrant."""

import json
import io
import uuid
from datetime import datetime, timezone

import pymysql
import tiktoken
from sentence_transformers import SentenceTransformer
from qdrant_client import QdrantClient
from qdrant_client.models import (
    Distance,
    VectorParams,
    PointStruct,
    Filter,
    FieldCondition,
    MatchValue,
)

from config import Config
from processors.embedding_config import EmbeddingConfig
from processors.storage import get_object

# Lazy-loaded globals
_model = None
_qdrant = None
_tokenizer = None


def get_embedding_model():
    global _model
    if _model is None:
        print(f"Loading embedding model: {EmbeddingConfig.EMBEDDING_MODEL}")
        _model = SentenceTransformer(EmbeddingConfig.EMBEDDING_MODEL)
    return _model


def get_qdrant_client():
    global _qdrant
    if _qdrant is None:
        _qdrant = QdrantClient(host=EmbeddingConfig.QDRANT_HOST, port=EmbeddingConfig.QDRANT_PORT)
        _ensure_collection()
    return _qdrant


def get_tokenizer():
    global _tokenizer
    if _tokenizer is None:
        _tokenizer = tiktoken.get_encoding("cl100k_base")
    return _tokenizer


def _ensure_collection():
    """Create Qdrant collection if it doesn't exist."""
    client = _qdrant
    collections = [c.name for c in client.get_collections().collections]
    if EmbeddingConfig.COLLECTION_NAME not in collections:
        client.create_collection(
            collection_name=EmbeddingConfig.COLLECTION_NAME,
            vectors_config=VectorParams(
                size=EmbeddingConfig.EMBEDDING_DIMENSION,
                distance=Distance.COSINE,
            ),
        )
        # Create payload indexes for filtering
        client.create_payload_index(EmbeddingConfig.COLLECTION_NAME, "organization_id", "integer")
        client.create_payload_index(EmbeddingConfig.COLLECTION_NAME, "workspace_id", "integer")
        client.create_payload_index(EmbeddingConfig.COLLECTION_NAME, "document_id", "keyword")
        print(f"Created Qdrant collection: {EmbeddingConfig.COLLECTION_NAME}")


def get_db_connection():
    return pymysql.connect(
        host=Config.MYSQL_HOST,
        port=Config.MYSQL_PORT,
        user=Config.MYSQL_USER,
        password=Config.MYSQL_PASSWORD,
        database=Config.MYSQL_DATABASE,
        cursorclass=pymysql.cursors.DictCursor,
    )


def chunk_text(text: str, page_number: int = 1) -> list[dict]:
    """Split text into chunks using recursive character splitting with token counting."""
    tokenizer = get_tokenizer()
    chunks = []
    separators = ["\n\n", "\n", ". ", " "]

    def _split(text_segment, char_offset, page_num):
        tokens = tokenizer.encode(text_segment)
        if len(tokens) <= EmbeddingConfig.CHUNK_SIZE:
            if len(tokens) >= EmbeddingConfig.MIN_CHUNK_SIZE:
                chunks.append({
                    "text": text_segment.strip(),
                    "char_start": char_offset,
                    "char_end": char_offset + len(text_segment),
                    "page_number": page_num,
                    "token_count": len(tokens),
                })
            return

        # Find best separator
        best_sep = " "
        for sep in separators:
            if sep in text_segment:
                best_sep = sep
                break

        parts = text_segment.split(best_sep)
        current_chunk = ""
        current_offset = char_offset

        for part in parts:
            candidate = current_chunk + best_sep + part if current_chunk else part
            candidate_tokens = tokenizer.encode(candidate)

            if len(candidate_tokens) > EmbeddingConfig.CHUNK_SIZE and current_chunk:
                chunk_tokens = tokenizer.encode(current_chunk)
                if len(chunk_tokens) >= EmbeddingConfig.MIN_CHUNK_SIZE:
                    chunks.append({
                        "text": current_chunk.strip(),
                        "char_start": current_offset,
                        "char_end": current_offset + len(current_chunk),
                        "page_number": page_num,
                        "token_count": len(chunk_tokens),
                    })
                # Overlap: keep the end of the current chunk
                overlap_text = current_chunk[-200:] if len(current_chunk) > 200 else ""
                current_offset = current_offset + len(current_chunk) - len(overlap_text)
                current_chunk = overlap_text + best_sep + part if overlap_text else part
            else:
                current_chunk = candidate

        # Remaining text
        if current_chunk:
            chunk_tokens = tokenizer.encode(current_chunk)
            if len(chunk_tokens) >= EmbeddingConfig.MIN_CHUNK_SIZE:
                chunks.append({
                    "text": current_chunk.strip(),
                    "char_start": current_offset,
                    "char_end": current_offset + len(current_chunk),
                    "page_number": page_num,
                    "token_count": len(chunk_tokens),
                })

    _split(text, 0, page_number)
    return chunks


def extract_text_from_file(bucket: str, key: str, mime_type: str) -> str:
    """Download file from PostgreSQL storage and extract text content."""
    content = get_object(bucket, key)

    if mime_type == "text/plain":
        return content.decode("utf-8", errors="replace")
    elif mime_type == "application/pdf":
        try:
            import fitz  # PyMuPDF
            doc = fitz.open(stream=content, filetype="pdf")
            text_parts = []
            for page in doc:
                text_parts.append(page.get_text())
            return "\n\n".join(text_parts)
        except ImportError:
            # Fallback: try pdfplumber
            import pdfplumber
            with pdfplumber.open(io.BytesIO(content)) as pdf:
                return "\n\n".join(page.extract_text() or "" for page in pdf.pages)
    elif mime_type in (
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "application/msword",
    ):
        try:
            from docx import Document
            doc = Document(io.BytesIO(content))
            return "\n\n".join(para.text for para in doc.paragraphs if para.text.strip())
        except ImportError:
            return ""
    elif mime_type in (
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "application/vnd.ms-powerpoint",
    ):
        try:
            from pptx import Presentation
            prs = Presentation(io.BytesIO(content))
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        texts.append(shape.text)
            return "\n\n".join(texts)
        except ImportError:
            return ""
    elif mime_type in (
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "application/vnd.ms-excel",
    ):
        try:
            import openpyxl
            wb = openpyxl.load_workbook(io.BytesIO(content), read_only=True)
            texts = []
            for ws in wb.worksheets:
                for row in ws.iter_rows(values_only=True):
                    row_text = " | ".join(str(cell) for cell in row if cell is not None)
                    if row_text.strip():
                        texts.append(row_text)
            return "\n".join(texts)
        except ImportError:
            return ""

    return ""


def process_embedding(job_data: dict):
    """Main embedding processor: extract text, chunk, embed, upsert to Qdrant."""
    file_id = job_data.get("fileId")
    file_uuid = job_data.get("fileUuid")
    org_id = job_data.get("organizationId")
    workspace_id = job_data.get("workspaceId")
    storage_bucket = job_data.get("storageBucket")
    storage_key = job_data.get("storageKey")
    mime_type = job_data.get("mimeType", "")
    file_name = job_data.get("fileName", "")
    job_uuid = job_data.get("jobUuid")

    print(f"[Embedding] Processing file {file_uuid} ({file_name}, mime={mime_type})")

    conn = get_db_connection()
    try:
        # Update job status to PROCESSING
        with conn.cursor() as cursor:
            cursor.execute(
                "UPDATE embedding_jobs SET status='PROCESSING', started_at=NOW() WHERE uuid=%s",
                (job_uuid,),
            )
            conn.commit()

        # Extract text
        text = extract_text_from_file(storage_bucket, storage_key, mime_type)

        if not text.strip():
            with conn.cursor() as cursor:
                cursor.execute(
                    "UPDATE embedding_jobs SET status='COMPLETED', chunk_count=0, completed_at=NOW() WHERE uuid=%s",
                    (job_uuid,),
                )
                conn.commit()
            print(f"[Embedding] No text extracted from {file_uuid}, marking completed with 0 chunks")
            return

        # Chunk text
        chunks = chunk_text(text)
        print(f"[Embedding] Generated {len(chunks)} chunks for {file_uuid}")

        # Generate embeddings in batches
        model = get_embedding_model()
        qdrant = get_qdrant_client()

        # Delete existing vectors for this document (idempotent re-indexing)
        qdrant.delete(
            collection_name=EmbeddingConfig.COLLECTION_NAME,
            points_selector=Filter(
                must=[FieldCondition(key="document_id", match=MatchValue(value=file_uuid))]
            ),
        )

        # Batch embed and upsert
        points = []
        for batch_start in range(0, len(chunks), EmbeddingConfig.BATCH_SIZE):
            batch = chunks[batch_start: batch_start + EmbeddingConfig.BATCH_SIZE]
            texts = [c["text"] for c in batch]
            embeddings = model.encode(texts).tolist()

            for i, (chunk, embedding) in enumerate(zip(batch, embeddings)):
                point_id = str(uuid.uuid4())
                points.append(
                    PointStruct(
                        id=point_id,
                        vector=embedding,
                        payload={
                            "document_id": file_uuid,
                            "organization_id": org_id,
                            "workspace_id": workspace_id,
                            "page_number": chunk["page_number"],
                            "char_start": chunk["char_start"],
                            "char_end": chunk["char_end"],
                            "chunk_text": chunk["text"],
                            "chunk_index": batch_start + i,
                            "file_name": file_name,
                            "mime_type": mime_type,
                            "embedded_at": datetime.now(timezone.utc).isoformat(),
                            "embedding_model": EmbeddingConfig.EMBEDDING_MODEL,
                        },
                    )
                )

        # Upsert in batches of 100 points
        for i in range(0, len(points), 100):
            qdrant.upsert(
                collection_name=EmbeddingConfig.COLLECTION_NAME,
                points=points[i: i + 100],
            )

        # Update job status to COMPLETED
        with conn.cursor() as cursor:
            cursor.execute(
                "UPDATE embedding_jobs SET status='COMPLETED', chunk_count=%s, completed_at=NOW() WHERE uuid=%s",
                (len(chunks), job_uuid),
            )
            conn.commit()

        print(f"[Embedding] Successfully embedded {len(chunks)} chunks for {file_uuid}")

    except Exception as e:
        # Update job status to FAILED with retry tracking
        error_msg = str(e)[:1000]
        try:
            with conn.cursor() as cursor:
                # Get current retry count
                cursor.execute(
                    "SELECT retry_count FROM embedding_jobs WHERE uuid=%s", (job_uuid,)
                )
                row = cursor.fetchone()
                current_retries = row["retry_count"] if row and row.get("retry_count") else 0

                if current_retries >= 2:
                    # Max 3 attempts (0, 1, 2) - move to DLQ
                    cursor.execute(
                        "UPDATE embedding_jobs SET status='DEAD_LETTER', error_message=%s, completed_at=NOW() WHERE uuid=%s",
                        (f"DLQ after {current_retries + 1} attempts: {error_msg}", job_uuid),
                    )
                    print(f"[Embedding] Job {job_uuid} moved to DLQ after {current_retries + 1} attempts: {error_msg}")
                else:
                    # Mark as FAILED with incremented retry count for reprocessing
                    cursor.execute(
                        "UPDATE embedding_jobs SET status='FAILED', retry_count=%s, error_message=%s WHERE uuid=%s",
                        (current_retries + 1, error_msg, job_uuid),
                    )
                    print(f"[Embedding] Job {job_uuid} failed (attempt {current_retries + 1}/3): {error_msg}")
                conn.commit()
        except Exception as db_err:
            print(f"[Embedding] Failed to update job status for {job_uuid}: {db_err}")
        raise
    finally:
        conn.close()
